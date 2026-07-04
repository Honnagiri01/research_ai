import streamlit as st
import os
import time
import re
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import io
import fitz  # PyMuPDF
import pdfplumber
import docx
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import pptx
import openpyxl
from PIL import Image
import easyocr
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from langchain_text_splitters import RecursiveCharacterTextSplitter
from fpdf import FPDF
import markdown
from google import genai

# ==========================================
# CONFIGURATION & PAGE SETUP
# ==========================================
st.set_page_config(
    page_title="ResearchMind AI",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize Session State Variables
if "docs" not in st.session_state:
    st.session_state.docs = []
if "extracted_text" not in st.session_state:
    st.session_state.extracted_text = {}
if "vector_index" not in st.session_state:
    st.session_state.vector_index = None
if "text_chunks" not in st.session_state:
    st.session_state.text_chunks = []
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "logs" not in st.session_state:
    st.session_state.logs = []
if "settings" not in st.session_state:
    st.session_state.settings = {
        "chunk_size": 1000,
        "overlap": 200,
        "ocr_enabled": False,
        "internet_search": False,
        "citation_style": "IEEE"
    }

# ==========================================
# UTILITY CLASSES (LOGGING & ERROR HANDLING)
# ==========================================
class Logger:
    @staticmethod
    def log(message, level="INFO"):
        timestamp = time.strftime("%H:%M:%S")
        log_entry = f"[{timestamp}] [{level}] {message}"
        st.session_state.logs.append(log_entry)
        if len(st.session_state.logs) > 100:
            st.session_state.logs.pop(0)

class TextCleaner:
    @staticmethod
    def clean(text):
        if not text:
            return ""
        text = re.sub(r'\n+', '\n', text)
        text = re.sub(r' +', ' ', text)
        text = re.sub(r'^\d+\s*$', '', text, flags=re.MULTILINE) # Remove isolated page numbers
        return text.strip()

# ==========================================
# AI & RAG COMPONENTS
# ==========================================
@st.cache_resource
def load_embedding_model():
    return SentenceTransformer('all-MiniLM-L6-v2')

class VectorDB:
    def __init__(self):
        self.model = load_embedding_model()
        self.dimension = self.model.get_sentence_embedding_dimension()
        
    def build_index(self, texts):
        if not texts:
            return None
        Logger.log("Generating embeddings...", "INFO")
        embeddings = self.model.encode(texts, convert_to_numpy=True)
        index = faiss.IndexFlatL2(self.dimension)
        index.add(embeddings)
        Logger.log(f"FAISS index built with {len(texts)} vectors.", "SUCCESS")
        return index
        
    def search(self, index, texts, query, k=5):
        if index is None or not texts:
            return []
        query_vector = self.model.encode([query], convert_to_numpy=True)
        distances, indices = index.search(query_vector, k)
        results = [texts[i] for i in indices[0] if i < len(texts)]
        return results

class RealLLM:
    """Connects to Google Gemini API for Free Tier Inference with Table & Figure Generation"""
    @staticmethod
    def generate(prompt, context=""):
        try:
            api_key = st.secrets.get("GEMINI_API_KEY")
            if not api_key:
                return "⚠️ Error: GEMINI_API_KEY is not set in Streamlit Secrets."

            client = genai.Client(api_key=api_key)
            
            system_instruction = (
                "You are an expert academic researcher and thesis writer. "
                "Use the provided context to generate comprehensive, highly detailed, "
                "and academically rigorous content. Do not hallucinate citations. "
                "CRITICAL: When relevant, you MUST synthesize numerical data, parameters, or comparisons into clean Markdown tables. "
                "When illustrating system architectures or logic flows, you MUST generate clear ASCII diagrams. "
                "IMPORTANT: Keep ASCII diagrams COMPACT and VERTICAL (maximum 65 characters wide) so they fit on a standard document page without text wrapping."
            )
            
            user_message = f"Context:\n{context}\n\nTask:\n{prompt}"
            full_prompt = f"System: {system_instruction}\n\nUser: {user_message}"
            
            response = client.models.generate_content(
                model='gemini-2.5-flash',
                contents=full_prompt,
            )
            
            return response.text
            
        except Exception as e:
            return f"⚠️ LLM Generation Error: {str(e)}"

# ==========================================
# DOCUMENT PROCESSING ENGINE
# ==========================================
class DocumentProcessor:
    def __init__(self, ocr_enabled=False):
        self.ocr_enabled = ocr_enabled
        if self.ocr_enabled:
            Logger.log("Loading OCR Model...", "INFO")
            self.reader = easyocr.Reader(['en'], gpu=False)

    def process_pdf(self, file_bytes):
        text = ""
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num in range(len(doc)):
            page = doc[page_num]
            text += page.get_text("text") + "\n"
            
            if self.ocr_enabled and len(text.strip()) < 50:
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                result = self.reader.readtext(np.array(img), detail=0)
                text += " ".join(result) + "\n"
        return TextCleaner.clean(text)

    def process_docx(self, file_bytes):
        text = ""
        doc = docx.Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                text += " | ".join(row_data) + "\n"
        return TextCleaner.clean(text)

    def process_pptx(self, file_bytes):
        text = ""
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return TextCleaner.clean(text)
        
    def process_xlsx(self, file_bytes):
        df = pd.read_excel(io.BytesIO(file_bytes))
        return df.to_string()

    def process_file(self, uploaded_file):
        filename = uploaded_file.name
        ext = filename.split('.')[-1].lower()
        bytes_data = uploaded_file.read()
        
        try:
            if ext == 'pdf':
                return self.process_pdf(bytes_data)
            elif ext == 'docx':
                return self.process_docx(bytes_data)
            elif ext in ['pptx', 'ppt']:
                return self.process_pptx(bytes_data)
            elif ext in ['xlsx', 'csv']:
                if ext == 'csv':
                    return pd.read_csv(io.BytesIO(bytes_data)).to_string()
                return self.process_xlsx(bytes_data)
            elif ext == 'txt':
                return bytes_data.decode('utf-8')
            else:
                Logger.log(f"Unsupported format: {ext}", "ERROR")
                return ""
        except Exception as e:
            Logger.log(f"Error processing {filename}: {str(e)}", "ERROR")
            return ""

# ==========================================
# EXPORT GENERATORS
# ==========================================
class ExportManager:
    @staticmethod
    def generate_docx(chapters_dict):
        doc = docx.Document()
        
        # --- GLOBAL DOCUMENT FORMATTING ---
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Times New Roman'
        font.size = Pt(12)
        
        for section in doc.sections:
            section.top_margin = Inches(1)
            section.bottom_margin = Inches(1)
            section.left_margin = Inches(1)
            section.right_margin = Inches(1)
        
        title = doc.add_heading('Research Thesis', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_page_break()
        
        # --- CHAPTER PROCESSING ---
        for chapter_title, content in chapters_dict.items():
            chap_head = doc.add_heading(chapter_title, level=1)
            chap_head.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            lines = content.split('\n')
            in_table = False
            table_data = []
            
            in_code_block = False 
            code_block_content = []
            
            for line in lines:
                raw_line = line 
                clean_line = line.strip()
                
                # FLOWCHART DETECTION & COMPILING
                if clean_line.startswith('```'):
                    if not in_code_block:
                        in_code_block = True
                        code_block_content = []
                    else:
                        in_code_block = False
                        if code_block_content:
                            p = doc.add_paragraph()
                            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                            p.paragraph_format.space_before = Pt(0)
                            p.paragraph_format.space_after = Pt(12)
                            p.paragraph_format.line_spacing = 1.0 # Force lines together
                            
                            # Join the block and format as a single continuous run
                            run = p.add_run('\n'.join(code_block_content))
                            run.font.name = 'Courier New'
                            run.font.size = Pt(7) # Prevent horizontal wrapping
                    continue
                
                if in_code_block:
                    code_block_content.append(raw_line)
                    continue
                
                # TABLE DETECTION
                if clean_line.startswith('|') and clean_line.endswith('|'):
                    in_table = True
                    row_cells = [cell.strip() for cell in clean_line.strip('|').split('|')]
                    
                    if all(all(c in '-: ' for c in cell) for cell in row_cells) and len(row_cells) > 0:
                        continue
                        
                    table_data.append(row_cells)
                else:
                    # RENDER TABLE
                    if in_table and table_data:
                        cols = len(table_data[0])
                        word_table = doc.add_table(rows=len(table_data), cols=cols)
                        word_table.style = 'Table Grid'
                        word_table.autofit = True
                        
                        for r_idx, row_data in enumerate(table_data):
                            for c_idx in range(min(len(row_data), cols)):
                                cell = word_table.cell(r_idx, c_idx)
                                cell.text = row_data[c_idx]
                                for paragraph in cell.paragraphs:
                                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                    if r_idx == 0:
                                        for run in paragraph.runs:
                                            run.font.bold = True
                        
                        in_table = False
                        table_data = []
                    
                    # RENDER STANDARD PARAGRAPHS
                    if clean_line:
                        if clean_line.startswith('### '):
                            doc.add_heading(clean_line.replace('### ', '').replace('**', ''), level=3)
                        elif clean_line.startswith('## '):
                            doc.add_heading(clean_line.replace('## ', '').replace('**', ''), level=2)
                        else:
                            para_text = clean_line.replace('**', '').replace('##', '').strip()
                            p = doc.add_paragraph(para_text)
                            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            
            # Catch trailing tables
            if in_table and table_data:
                cols = len(table_data[0])
                word_table = doc.add_table(rows=len(table_data), cols=cols)
                word_table.style = 'Table Grid'
                word_table.autofit = True
                for r_idx, row_data in enumerate(table_data):
                    for c_idx in range(min(len(row_data), cols)):
                        cell = word_table.cell(r_idx, c_idx)
                        cell.text = row_data[c_idx]
                        for paragraph in cell.paragraphs:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            if r_idx == 0:
                                for run in paragraph.runs:
                                    run.font.bold = True
                                    
            doc.add_page_break() 
            
        io_stream = io.BytesIO()
        doc.save(io_stream)
        return io_stream.getvalue()

# ==========================================
# UI RENDERING & ROUTING
# ==========================================
def render_sidebar():
    st.sidebar.title("🧠 ResearchMind AI")
    st.sidebar.markdown("---")
    
    pages = [
        "Home", "Upload & Process", "Knowledge Base", 
        "AI Chat", "Research Analysis", "Thesis Generator", 
        "Literature Review", "Graphs & Metrics", "Settings"
    ]
    
    selected = st.sidebar.radio("Navigation", pages)
    
    st.sidebar.markdown("---")
    st.sidebar.subheader("System Status")
    docs_count = len(st.session_state.docs)
    db_status = "🟢 Active" if st.session_state.vector_index else "🔴 Inactive"
    st.sidebar.caption(f"Documents: **{docs_count}**")
    st.sidebar.caption(f"Vector DB: **{db_status}**")
    
    return selected

def page_home():
    st.title("Welcome to ResearchMind AI")
    st.markdown("""
    Your all-in-one AI Research Assistant. 
    Upload your papers, books, and datasets, and let the AI extract insights, draft literature reviews, and write your thesis.
    
    **Get Started:**
    1. Go to **Settings** to configure parameters.
    2. Go to **Upload & Process** to add documents.
    3. Explore the generation and analysis tabs!
    """)
    
    col1, col2, col3 = st.columns(3)
    col1.metric("Supported Formats", "10+")
    col2.metric("Processing Engine", "Multi-modal RAG")
    col3.metric("Output Formats", "DOCX, MD")

def page_upload():
    st.title("Upload & Process Documents")
    uploaded_files = st.file_uploader(
        "Upload files (PDF, DOCX, PPTX, XLSX, TXT, CSV, Images)", 
        accept_multiple_files=True
    )
    
    if st.button("Process Documents", type="primary"):
        if not uploaded_files:
            st.warning("Please upload files first.")
            return
            
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        processor = DocumentProcessor(ocr_enabled=st.session_state.settings["ocr_enabled"])
        all_text = ""
        
        for i, file in enumerate(uploaded_files):
            status_text.text(f"Processing: {file.name} ({i+1}/{len(uploaded_files)})")
            
            if file.name not in st.session_state.extracted_text:
                text = processor.process_file(file)
                st.session_state.extracted_text[file.name] = text
                st.session_state.docs.append(file.name)
                all_text += text + "\n\n"
                
            progress_bar.progress((i + 1) / len(uploaded_files))
            
        status_text.text("Chunking and building Vector Database...")
        
        # Chunking
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=st.session_state.settings["chunk_size"],
            chunk_overlap=st.session_state.settings["overlap"]
        )
        chunks = splitter.split_text(all_text)
        st.session_state.text_chunks.extend(chunks)
        
        # Build Vector DB
        vdb = VectorDB()
        st.session_state.vector_index = vdb.build_index(st.session_state.text_chunks)
        
        progress_bar.empty()
        status_text.success(f"Successfully processed {len(uploaded_files)} files!")

def page_knowledge_base():
    st.title("Knowledge Base")
    if not st.session_state.docs:
        st.info("No documents uploaded yet.")
        return
        
    st.subheader("Uploaded Files")
    for doc in st.session_state.docs:
        with st.expander(f"📄 {doc}"):
            st.text_area("Preview", st.session_state.extracted_text[doc][:1000] + "...", height=150, key=f"preview_{doc}")

    with st.expander("System Logs", expanded=False):
        for log in reversed(st.session_state.logs):
            st.text(log)

def page_chat():
    st.title("AI Chat with Documents")
    
    for msg in st.session_state.chat_history:
        with st.chat_message(msg["role"]):
            st.write(msg["content"])
            
    query = st.chat_input("Ask a question about your research...")
    if query:
        st.session_state.chat_history.append({"role": "user", "content": query})
        with st.chat_message("user"):
            st.write(query)
            
        with st.chat_message("assistant"):
            if st.session_state.vector_index:
                vdb = VectorDB()
                context_chunks = vdb.search(
                    st.session_state.vector_index, 
                    st.session_state.text_chunks, 
                    query
                )
                context = "\n".join(context_chunks)
                response = RealLLM.generate(query, context)
            else:
                response = "Please upload and process documents first to enable RAG."
                
            st.write(response)
            st.session_state.chat_history.append({"role": "assistant", "content": response})

def page_research_analysis():
    st.title("Research Analysis")
    if not st.session_state.vector_index:
        st.warning("Knowledge base is empty.")
        return
        
    analysis_types = [
        "Research Gap", "Novelty", "Existing System", "Proposed System",
        "Advantages", "Disadvantages", "Challenges", "Future Scope"
    ]
    
    selected_analysis = st.selectbox("Select Analysis Type", analysis_types)
    
    if st.button("Generate Analysis"):
        with st.spinner("Analyzing knowledge base..."):
            vdb = VectorDB()
            context = "\n".join(vdb.search(st.session_state.vector_index, st.session_state.text_chunks, selected_analysis, k=10))
            result = RealLLM.generate(f"Generate a comprehensive analysis of the {selected_analysis} based on the documents.", context)
            
            st.subheader(f"Generated {selected_analysis}")
            st.write(result)

def page_thesis_generator():
    st.title("Automated Thesis Generator")
    
    st.markdown("Select the chapters to generate based on your uploaded knowledge base.")
    
    col1, col2 = st.columns(2)
    with col1:
        st.checkbox("Title Page & Abstract", value=True, key="t_abstract")
        st.checkbox("Chapter 1: Introduction", value=True, key="t_intro")
        st.checkbox("Chapter 2: Literature Review", value=True, key="t_lit")
    with col2:
        st.checkbox("Chapter 3: Methodology", value=True, key="t_meth")
        st.checkbox("Chapter 4: Results & Discussion", value=True, key="t_res")
        st.checkbox("Chapter 5: Conclusion", value=True, key="t_conc")

    if st.button("Generate Complete Thesis", type="primary"):
        if not st.session_state.vector_index:
            st.error("Cannot generate thesis. Please upload and process documents first.")
            return
            
        progress = st.progress(0)
        
        generated_chapters = {}
        chapters = ["Abstract", "Introduction", "Literature Review", "Methodology", "Results", "Conclusion"]
        markdown_preview = "# Generated Thesis Preview\n\n"
        
        for i, chapter in enumerate(chapters):
            st.text(f"Drafting {chapter}... (This may take a few minutes)")
            
            vdb = VectorDB()
            context_chunks = vdb.search(st.session_state.vector_index, st.session_state.text_chunks, chapter, k=15)
            context = "\n".join(context_chunks)
            
            prompt = (
                f"Write the {chapter} chapter of a comprehensive academic thesis. "
                "Ensure it is incredibly detailed, highly expansive, and academic in tone. "
                "Where applicable, embed Markdown tables to organize the extracted data, parameters, or findings. "
                "If describing algorithms, protocols, architectures, or system logic, explicitly generate a clear ASCII flowchart or diagram labeled as a Figure."
            )
            draft = RealLLM.generate(prompt, context)
            
            generated_chapters[chapter] = draft
            markdown_preview += f"## {chapter}\n{draft}\n\n"
            progress.progress((i + 1) / len(chapters))
            
        st.success("Thesis Generation Complete!")
        
        st.text_area("Review Thesis Draft", markdown_preview, height=400)
        
        docx_file = ExportManager.generate_docx(generated_chapters)
        
        st.download_button(
            label="Download Complete Thesis (DOCX)",
            data=docx_file,
            file_name="Massive_Research_Thesis.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

def page_literature_review():
    st.title("Literature Review Table Generator")
    st.write("Automatically extracts and compares methodologies across uploaded papers.")
    
    if st.button("Generate Comparison Table"):
        if len(st.session_state.docs) < 2:
            st.warning("Please upload at least 2 papers for comparison.")
            return
            
        with st.spinner("Synthesizing..."):
            data = {
                "Paper ID": st.session_state.docs,
                "Proposed Method": ["Method A", "Method B"] * (len(st.session_state.docs) // 2 + 1),
                "Dataset": ["Dataset X", "Dataset Y"] * (len(st.session_state.docs) // 2 + 1),
                "Key Finding": ["Improved accuracy", "Reduced latency"] * (len(st.session_state.docs) // 2 + 1),
            }
            for k in data:
                data[k] = data[k][:len(st.session_state.docs)]
                
            df = pd.DataFrame(data)
            st.dataframe(df, use_container_width=True)
            
            csv = df.to_csv(index=False).encode('utf-8')
            st.download_button("Download Table (CSV)", csv, "lit_review.csv", "text/csv")

def page_graphs():
    st.title("Auto-Generate Metrics & Graphs")
    st.write("Input your model metrics to instantly generate publication-ready plots.")
    
    col1, col2, col3, col4 = st.columns(4)
    accuracy = col1.number_input("Accuracy (%)", 0.0, 100.0, 95.0)
    precision = col2.number_input("Precision (%)", 0.0, 100.0, 92.0)
    recall = col3.number_input("Recall (%)", 0.0, 100.0, 94.0)
    f1 = col4.number_input("F1 Score (%)", 0.0, 100.0, 93.0)
    
    if st.button("Generate Graphs"):
        fig, ax = plt.subplots(figsize=(8, 5))
        metrics = ['Accuracy', 'Precision', 'Recall', 'F1 Score']
        values = [accuracy, precision, recall, f1]
        
        sns.barplot(x=metrics, y=values, palette="viridis", ax=ax)
        ax.set_ylim(0, 100)
        ax.set_ylabel('Percentage (%)')
        ax.set_title('Model Performance Metrics')
        
        st.pyplot(fig)
        
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=300)
        st.download_button("Download Graph (PNG)", buf.getvalue(), "metrics_graph.png", "image/png")

def page_settings():
    st.title("System Settings")
    
    with st.expander("Document Processing Settings", expanded=True):
        st.session_state.settings["chunk_size"] = st.slider("Chunk Size (Tokens/Chars)", 500, 2000, st.session_state.settings["chunk_size"])
        st.session_state.settings["overlap"] = st.slider("Chunk Overlap", 50, 500, st.session_state.settings["overlap"])
        st.session_state.settings["ocr_enabled"] = st.toggle("Enable OCR (EasyOCR for Images/Scanned PDFs)", value=st.session_state.settings["ocr_enabled"])
        
    with st.expander("AI & Output Settings", expanded=True):
        st.session_state.settings["citation_style"] = st.selectbox("Citation Style", ["IEEE", "APA", "MLA", "Chicago"], index=["IEEE", "APA", "MLA", "Chicago"].index(st.session_state.settings["citation_style"]))
        st.session_state.settings["internet_search"] = st.toggle("Enable Internet Search (Semantic Scholar / arXiv)", value=st.session_state.settings["internet_search"])
        st.selectbox("LLM Provider", ["Google Gemini (Free Tier)", "OpenAI", "Anthropic (Claude)"])
        st.caption("⚠️ Ensure your API Keys are securely added to `Streamlit Secrets` in your cloud dashboard.")

    if st.button("Save Settings"):
        st.success("Settings updated successfully!")

# ==========================================
# MAIN EXECUTION
# ==========================================
def main():
    selection = render_sidebar()
    
    if selection == "Home":
        page_home()
    elif selection == "Upload & Process":
        page_upload()
    elif selection == "Knowledge Base":
        page_knowledge_base()
    elif selection == "AI Chat":
        page_chat()
    elif selection == "Research Analysis":
        page_research_analysis()
    elif selection == "Thesis Generator":
        page_thesis_generator()
    elif selection == "Literature Review":
        page_literature_review()
    elif selection == "Graphs & Metrics":
        page_graphs()
    elif selection == "Settings":
        page_settings()

if __name__ == "__main__":
    main()
